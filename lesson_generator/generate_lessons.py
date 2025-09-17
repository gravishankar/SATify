
import os
import sys
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas

def parse_markdown(md_path):
    """Parse lesson.md into sections dictionary."""
    sections = {}
    current_section = None
    with open(md_path, "r") as f:
        for line in f:
            line = line.strip()
            if line.startswith("## "):  # section header
                current_section = line.replace("## ", "").strip()
                sections[current_section] = ""
            elif line.startswith("# "):  # skip top-level title
                continue
            elif current_section:
                sections[current_section] += line + "\n"
    return sections

def generate_lesson(md_path):
    """Generate PPTX, HTML, PDF from lesson.md"""
    skill_dir = os.path.dirname(md_path)
    skill_name = os.path.basename(skill_dir)
    sections = parse_markdown(md_path)

    # -------- PPTX --------
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0]).shapes.title.text = f"SAT Skill: {skill_name.capitalize()}"

    def add_slide(prs, title, content):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = content

    for title, content in sections.items():
        add_slide(prs, title, content.strip())

    pptx_path = os.path.join(skill_dir, "deck.pptx")
    prs.save(pptx_path)

    # -------- HTML (reveal.js) --------
    html_content = f"""
    <!doctype html>
    <html lang="en">
    <head>
      <meta charset="utf-8">
      <title>{skill_name.capitalize()} - SAT Skill</title>
      <link rel="stylesheet" href="https://unpkg.com/reveal.js/dist/reveal.css">
      <link rel="stylesheet" href="https://unpkg.com/reveal.js/dist/theme/white.css">
    </head>
    <body>
      <div class="reveal">
        <div class="slides">
          <section><h2>SAT Skill: {skill_name.capitalize()}</h2></section>
    """

    for title, content in sections.items():
        html_content += f"<section><h3>{title}</h3><pre>{content.strip()}</pre></section>"

    html_content += """
        </div>
      </div>
      <script src="https://unpkg.com/reveal.js/dist/reveal.js"></script>
      <script>Reveal.initialize();</script>
    </body>
    </html>
    """

    html_path = os.path.join(skill_dir, "deck.html")
    with open(html_path, "w") as f:
        f.write(html_content)

    # -------- PDF --------
    pdf_path = os.path.join(skill_dir, "strategy.pdf")
    c = canvas.Canvas(pdf_path, pagesize=letter)
    width, height = letter

    y = height - 50
    c.setFont("Helvetica-Bold", 14)
    c.drawCentredString(width/2, y, f"SAT Skill: {skill_name.capitalize()}")

    y -= 40
    for title, content in sections.items():
        c.setFont("Helvetica-Bold", 12)
        c.drawString(50, y, title)
        y -= 20
        c.setFont("Helvetica", 10)
        text_obj = c.beginText(50, y)
        for line in content.strip().split("\n"):
            text_obj.textLine(line)
            y -= 12
            if y < 100:
                c.drawText(text_obj)
                c.showPage()
                y = height - 50
                text_obj = c.beginText(50, y)
        c.drawText(text_obj)
        y -= 20

    c.save()

    print("Generated files:")
    print(" -", pptx_path)
    print(" -", html_path)
    print(" -", pdf_path)

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python generate_lessons.py path/to/lesson.md")
        sys.exit(1)
    md_path = sys.argv[1]
    if not os.path.exists(md_path):
        print("File not found:", md_path)
        sys.exit(1)
    generate_lesson(md_path)
