
    import os
    import sys
    from pptx import Presentation
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas

    def parse_markdown(md_path):
        """Parse lesson.md into sections dictionary."""
        sections = {}
        current_section = None
        with open(md_path, "r") as f:
            for line in f:
                line = line.strip()
                if line.startswith("## "):  # section header
                    current_section = line.replace("## ", "").strip()
                    sections[current_section] = ""
                elif line.startswith("# "):  # skip top-level title
                    continue
                elif current_section:
                    sections[current_section] += line + "\n"
        return sections

    def update_index(base_dir):
        """Rebuild index.html and README.md for all skills in base_dir."""
        index_path = os.path.join(base_dir, "index.html")
        readme_path = os.path.join(base_dir, "README.md")

        skills = []
        for folder in sorted(os.listdir(base_dir)):
            skill_dir = os.path.join(base_dir, folder)
            if os.path.isdir(skill_dir):
                if all(os.path.exists(os.path.join(skill_dir, f)) for f in ["deck.html", "deck.pptx", "strategy.pdf"]):
                    skills.append((folder.replace("-", " ").title(), folder))

        # Build index.html
        html = """<!doctype html>
<html lang="en">
<head>
  <meta charset="utf-8">
  <title>SAT Practice Pro – Skills Hub</title>
  <style>
    body { font-family: Arial, sans-serif; margin: 40px; background: #f9f9f9; color: #222; }
    h1 { color: #004080; }
    .skill-card { background: white; border: 1px solid #ddd; border-radius: 8px; padding: 20px; margin: 15px 0; }
    .skill-card h2 { margin-top: 0; }
    .links a { margin-right: 15px; text-decoration: none; color: #0066cc; }
    .links a:hover { text-decoration: underline; }
  </style>
</head>
<body>
  <h1>SAT Practice Pro – Skills Hub</h1>
  <p>Select a skill below to view its lesson materials:</p>
        """
        for title, folder in skills:
            html += f"""
    <div class="skill-card">
      <h2>{title}</h2>
      <div class="links">
        🌐 <a href="./{folder}/deck.html">Web Deck</a>
        📑 <a href="./{folder}/deck.pptx">PowerPoint</a>
        📄 <a href="./{folder}/strategy.pdf">PDF Strategy Sheet</a>
      </div>
    </div>"""
        html += "</body></html>"

        with open(index_path, "w") as f:
            f.write(html)

        # Build README.md
        md = "# SAT Practice Pro – Skills Index\n\n"
        for title, folder in skills:
            md += f"## {title}\n"
            md += f"- 🌐 [Web Deck](./{folder}/deck.html)\n"
            md += f"- 📑 [PowerPoint](./{folder}/deck.pptx)\n"
            md += f"- 📄 [PDF Strategy Sheet](./{folder}/strategy.pdf)\n\n"

        with open(readme_path, "w") as f:
            f.write(md)

    def generate_lesson(md_path):
        """Generate PPTX, HTML, PDF from lesson.md and update index."""
        skill_dir = os.path.dirname(md_path)
        base_dir = os.path.dirname(skill_dir)
        skill_name = os.path.basename(skill_dir)
        sections = parse_markdown(md_path)

        # -------- PPTX --------
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0]).shapes.title.text = f"SAT Skill: {skill_name.replace('-', ' ').title()}"

        def add_slide(prs, title, content):
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title
            slide.placeholders[1].text = content

        for title, content in sections.items():
            add_slide(prs, title, content.strip())

        pptx_path = os.path.join(skill_dir, "deck.pptx")
        prs.save(pptx_path)

        # -------- HTML (reveal.js) --------
        html_content = f"""
        <!doctype html>
        <html lang="en">
        <head>
          <meta charset="utf-8">
          <title>{skill_name.replace('-', ' ').title()} - SAT Skill</title>
          <link rel="stylesheet" href="https://unpkg.com/reveal.js/dist/reveal.css">
          <link rel="stylesheet" href="https://unpkg.com/reveal.js/dist/theme/white.css">
        </head>
        <body>
          <div class="reveal">
            <div class="slides">
              <section><h2>SAT Skill: {skill_name.replace('-', ' ').title()}</h2></section>
        """

        for title, content in sections.items():
            html_content += f"<section><h3>{title}</h3><pre>{content.strip()}</pre></section>"

        html_content += """
            </div>
          </div>
          <script src="https://unpkg.com/reveal.js/dist/reveal.js"></script>
          <script>Reveal.initialize();</script>
        </body>
        </html>
        """

        html_path = os.path.join(skill_dir, "deck.html")
        with open(html_path, "w") as f:
            f.write(html_content)

        # -------- PDF --------
        pdf_path = os.path.join(skill_dir, "strategy.pdf")
        c = canvas.Canvas(pdf_path, pagesize=letter)
        width, height = letter

        y = height - 50
        c.setFont("Helvetica-Bold", 14)
        c.drawCentredString(width/2, y, f"SAT Skill: {skill_name.replace('-', ' ').title()}")

        y -= 40
        for title, content in sections.items():
            c.setFont("Helvetica-Bold", 12)
            c.drawString(50, y, title)
            y -= 20
            c.setFont("Helvetica", 10)
            text_obj = c.beginText(50, y)
            for line in content.strip().split("\n"):
                text_obj.textLine(line)
                y -= 12
                if y < 100:
                    c.drawText(text_obj)
                    c.showPage()
                    y = height - 50
                    text_obj = c.beginText(50, y)
            c.drawText(text_obj)
            y -= 20

        c.save()

        # -------- Update Index --------
        update_index(base_dir)

        print("Generated files:")
        print(" -", pptx_path)
        print(" -", html_path)
        print(" -", pdf_path)
        print("Updated index.html and README.md in", base_dir)

    if __name__ == "__main__":
        if len(sys.argv) < 2:
            print("Usage: python generate_lessons.py path/to/lesson.md")
            sys.exit(1)
        md_path = sys.argv[1]
        if not os.path.exists(md_path):
            print("File not found:", md_path)
            sys.exit(1)
        generate_lesson(md_path)
